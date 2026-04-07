"""Generate a clean 5-slide presentation with properly centered text."""
import qrcode
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Palette from the app UI ---
BG        = RGBColor(0x20, 0x20, 0x20)
GREEN_DK  = RGBColor(0x28, 0x5A, 0x48)
GREEN     = RGBColor(0x40, 0x8A, 0x71)
GREEN_LT  = RGBColor(0x52, 0xB3, 0x93)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xCC, 0xCC, 0xCC)
CARD      = RGBColor(0x2A, 0x2A, 0x2A)
RED       = RGBColor(0xEF, 0x44, 0x44)

W = Inches(13.333)
H = Inches(7.5)
CX = W // 2  # center x

GITHUB_URL = "https://github.com/LimpingCoronation/se-toolkit-hackathon"
DEPLOY_URL = "https://your-deployed-url-here.com"


def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG


def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s


def card(slide, l, t, w, h, fill, title, body_lines,
         title_color=GREEN, body_color=GRAY, title_sz=18, body_sz=13):
    """Create a card shape with vertically & horizontally centered text."""
    shape = rrect(slide, l, t, w, h, fill)
    shape.text_frame.word_wrap = True
    shape.text_frame.margin_left = Inches(0.3)
    shape.text_frame.margin_right = Inches(0.3)
    shape.text_frame.margin_top = Inches(0.15)
    shape.text_frame.margin_bottom = Inches(0.15)
    tf = shape.text_frame
    # Vertical center
    shape.text_frame.auto_size = None
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_sz)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(6)

    for line in body_lines:
        np = tf.add_paragraph()
        np.text = line
        np.font.size = Pt(body_sz)
        np.font.color.rgb = body_color
        np.font.name = "Calibri"
        np.alignment = PP_ALIGN.CENTER
        np.space_after = Pt(3)

    return shape


def hdr(slide, title_text):
    """Green header bar with centered title."""
    rect(slide, Inches(0), Inches(0), W, Inches(1.0), GREEN_DK)
    tb = slide.shapes.add_textbox(Inches(0), Inches(0.15), W, Inches(0.7))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER


def para(tf, text, sz=14, bold=False, color=GRAY, align=PP_ALIGN.CENTER,
          before=Pt(0), after=Pt(0), italic=False):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.font.italic = italic
    p.alignment = align
    p.space_before = before
    p.space_after = after
    return p


def qr(url, path):
    code = qrcode.QRCode(version=1, box_size=10, border=2)
    code.add_data(url); code.make(fit=True)
    img = code.make_image(fill_color="#408A71", back_color="#202020").convert("RGB")
    img.save(path)


# ================================================================
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ======================== SLIDE 1 — TITLE ========================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

cw = Inches(9)
ch = Inches(4.0)
cl = (W - cw) // 2
ct = (H - ch) // 2 - Inches(0.2)

rrect(s, cl, ct, cw, ch, GREEN_DK)

# Accent line
rect(s, CX - Inches(1), ct + Inches(0.3), Inches(2), Pt(4), GREEN)

tb = s.shapes.add_textbox(cl + Inches(0.5), ct + Inches(0.7), cw - Inches(1), Inches(0.9))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.text = "Server Monitor"
p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE
p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER

tb2 = s.shapes.add_textbox(cl + Inches(0.5), ct + Inches(1.6), cw - Inches(1), Inches(0.5))
tb2.text_frame.word_wrap = True
p2 = tb2.text_frame.paragraphs[0]
p2.text = "Real-Time Server Health Dashboard"
p2.font.size = Pt(19); p2.font.color.rgb = GRAY
p2.font.name = "Calibri"; p2.alignment = PP_ALIGN.CENTER

# Divider
rect(s, cl + Inches(2), ct + Inches(2.25), cw - Inches(4), Pt(2), GREEN)

tb3 = s.shapes.add_textbox(cl + Inches(0.5), ct + Inches(2.5), cw - Inches(1), Inches(1.2))
tf3 = tb3.text_frame; tf3.word_wrap = True
para(tf3, "Your Name", 17, True, WHITE, before=Pt(0), after=Pt(4))
para(tf3, "your.email@university.edu", 14, False, GRAY, after=Pt(4))
para(tf3, "Group: XXXXX", 14, False, GRAY)

# Dot decoration
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cl + Inches(0.3), ct + Inches(0.78), Inches(0.12), Inches(0.12))
dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.fill.background()

# ======================== SLIDE 2 — CONTEXT ========================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
hdr(s, "Context")

cw2 = Inches(3.8)
ch2 = Inches(1.7)
gap = Inches(0.45)
total = cw2 * 3 + gap * 2
sx = (W - total) // 2
sy = Inches(1.5)

card(s, sx, sy, cw2, ch2, CARD,
     "👤  End User",
     ["System administrators and DevOps engineers",
      "who monitor multiple server endpoints"],
     title_sz=17, body_sz=13)

card(s, sx + cw2 + gap, sy, cw2, ch2, CARD,
     "⚠️  Problem",
     ["Manually checking server health is slow",
      "No centralized lightweight monitoring tool"],
     title_color=RED, title_sz=17, body_sz=13)

card(s, sx + (cw2 + gap) * 2, sy, cw2, ch2, CARD,
     "💡  Product Idea",
     ["Real-time dashboard tracking server health",
      "via WebSocket with instant alerts"],
     title_color=GREEN_LT, title_sz=17, body_sz=13)

# Quote box
qw = Inches(10)
ql = (W - qw) // 2
qb = rrect(s, ql, Inches(3.9), qw, Inches(1.3), GREEN_DK)
qb.text_frame.word_wrap = True
qb.text_frame.margin_left = Inches(0.6)
qb.text_frame.margin_right = Inches(0.6)
qb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = qb.text_frame.paragraphs[0]
p.text = '"Monitor all your servers in real time\nfrom one simple, beautiful dashboard."'
p.font.size = Pt(20); p.font.italic = True; p.font.color.rgb = WHITE
p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER

# ======================== SLIDE 3 — IMPLEMENTATION ========================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
hdr(s, "Implementation")

# Left: How we built it
lw = Inches(5.8)
lh = Inches(5.2)
lx = Inches(0.6)
ly = Inches(1.3)

lc = rrect(s, lx, ly, lw, lh, CARD)
lc.text_frame.word_wrap = True
lc.text_frame.margin_left = Inches(0.3)
lc.text_frame.margin_right = Inches(0.3)
lc.text_frame.margin_top = Inches(0.25)
lc.text_frame.margin_bottom = Inches(0.25)
tf = lc.text_frame
para(tf, "🛠  How We Built It", 19, True, GREEN, before=Pt(0), after=Pt(8))

items = [
    "Backend: FastAPI (Python 3.12+) with async",
    "Database: PostgreSQL + SQLAlchemy + Alembic",
    "Real-time: WebSocket + Redis Pub/Sub",
    "Auth: JWT token-based authentication",
    "Frontend: Vanilla HTML/CSS/JS, dark UI",
    "Containerized: Docker + docker-compose",
]
for it in items:
    para(tf, f"▹  {it}", 13, False, GRAY, before=Pt(0), after=Pt(5))

# Right top: V1 vs V2
rw = Inches(6.2)
rh = Inches(2.4)
rx = Inches(6.7)
ry = Inches(1.3)

rc = rrect(s, rx, ry, rw, rh, CARD)
rc.text_frame.word_wrap = True
rc.text_frame.margin_left = Inches(0.3)
rc.text_frame.margin_right = Inches(0.3)
rc.text_frame.margin_top = Inches(0.2)
rc.text_frame.margin_bottom = Inches(0.2)
tf2 = rc.text_frame
para(tf2, "📦  Version 1 → Version 2", 19, True, GREEN, before=Pt(0), after=Pt(6))

para(tf2, "Version 1:", 14, True, GREEN_LT, PP_ALIGN.LEFT, after=Pt(2))
for v in ["User authentication", "Server list CRUD", "Basic UI"]:
    para(tf2, f"  ▹  {v}", 12, False, GRAY, PP_ALIGN.LEFT, after=Pt(1))

para(tf2, "Version 2:", 14, True, GREEN_LT, PP_ALIGN.LEFT, before=Pt(5), after=Pt(2))
for v in ["WebSocket real-time monitoring", "Redis Pub/Sub live updates",
          "Start/stop controls", "Live status indicators"]:
    para(tf2, f"  ▹  {v}", 12, False, GRAY, PP_ALIGN.LEFT, after=Pt(1))

# Right bottom: TA Feedback
rfh = Inches(2.5)
rfy = ry + rh + Inches(0.3)

rf = rrect(s, rx, rfy, rw, rfh, CARD)
rf.text_frame.word_wrap = True
rf.text_frame.margin_left = Inches(0.3)
rf.text_frame.margin_right = Inches(0.3)
rf.text_frame.margin_top = Inches(0.2)
rf.text_frame.margin_bottom = Inches(0.2)
tf3 = rf.text_frame
para(tf3, "✅  TA Feedback Addressed", 19, True, GREEN, before=Pt(0), after=Pt(6))

ta = [
    "Error handling & HTTPException responses",
    "Clean code structure with routers & schemas",
    "JWT validation on WebSocket connections",
    "Proper DB session management",
    "CORS middleware for frontend-backend",
]
for it in ta:
    para(tf3, f"▹  {it}", 12, False, GRAY, before=Pt(0), after=Pt(3))

# ======================== SLIDE 4 — DEMO ========================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
hdr(s, "Demo")

# Video frame — centered
vw = Inches(10)
vh = Inches(4.0)
vl = (W - vw) // 2
vt = Inches(1.5)

rrect(s, vl, vt, vw, vh, CARD)

# Screen area
sw2 = Inches(9.4)
sh2 = Inches(3.4)
sl2 = (W - sw2) // 2
st2 = vt + Inches(0.3)

rect(s, sl2, st2, sw2, sh2, RGBColor(0x18, 0x18, 0x18))

# 4 corner labels
lbl_data = [
    (Inches(0.3), Inches(0.2), "Authentication Flow"),
    (sw2 - Inches(3.2), Inches(0.2), "Real-Time Monitoring"),
    (Inches(0.3), sh2 - Inches(0.5), "Server Management"),
    (sw2 - Inches(2.8), sh2 - Inches(0.5), "Live Status"),
]
for ox, oy, txt in lbl_data:
    tb = s.shapes.add_textbox(sl2 + ox, st2 + oy, Inches(3.5), Inches(0.4))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = txt; p.font.size = Pt(11); p.font.bold = True
    p.font.color.rgb = GREEN; p.font.name = "Calibri"; p.alignment = PP_ALIGN.LEFT

# Play button — exact center
pb = Inches(1.1)
px = sl2 + (sw2 - pb) // 2
py = st2 + (sh2 - pb) // 2
circ = s.shapes.add_shape(MSO_SHAPE.OVAL, px, py, pb, pb)
circ.fill.solid(); circ.fill.fore_color.rgb = GREEN; circ.line.fill.background()

tb = s.shapes.add_textbox(px, py + Inches(0.2), pb, Inches(0.7))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.text = "▶"; p.font.size = Pt(34); p.font.bold = True
p.font.color.rgb = WHITE; p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER

# Description — centered
tb = s.shapes.add_textbox(vl, vt + vh + Inches(0.2), vw, Inches(0.4))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.text = "Pre-recorded video demonstration of Version 2 with voice-over (max 2 min)"
p.font.size = Pt(14); p.font.color.rgb = GRAY; p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Important note — centered
nw = Inches(7)
nl = (W - nw) // 2
nb = rrect(s, nl, Inches(6.55), nw, Inches(0.5), GREEN_DK)
nb.text_frame.word_wrap = True
nb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = nb.text_frame.paragraphs[0]
p.text = "★  Most important part of the presentation"
p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE
p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER

# ======================== SLIDE 5 — LINKS ========================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
hdr(s, "Links")

cw5 = Inches(5.2)
ch5 = Inches(2.6)
gap5 = Inches(0.7)
total5 = cw5 * 2 + gap5
sx5 = (W - total5) // 2
sy5 = Inches(1.5)

# GitHub card
gc = rrect(s, sx5, sy5, cw5, ch5, CARD)
gc.text_frame.word_wrap = True
gc.text_frame.margin_left = Inches(0.35)
gc.text_frame.margin_right = Inches(0.35)
gc.text_frame.margin_top = Inches(0.25)
gc.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
tf = gc.text_frame
para(tf, "📂  GitHub Repository", 18, True, GREEN, before=Pt(0), after=Pt(6))
para(tf, GITHUB_URL, 11, False, GRAY, after=Pt(4))
para(tf, "LimpingCoronation/se-toolkit-hackathon", 13, True, WHITE)

# Deploy card
dx5 = sx5 + cw5 + gap5
dc = rrect(s, dx5, sy5, cw5, ch5, CARD)
dc.text_frame.word_wrap = True
dc.text_frame.margin_left = Inches(0.35)
dc.text_frame.margin_right = Inches(0.35)
dc.text_frame.margin_top = Inches(0.25)
dc.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
tf2 = dc.text_frame
para(tf2, "🚀  Deployed Product", 18, True, GREEN, before=Pt(0), after=Pt(6))
para(tf2, DEPLOY_URL, 11, False, GRAY, after=Pt(4))
para(tf2, "(Latest Version — V2)", 13, True, WHITE)

# QR codes
qr_gh = "/tmp/qr_gh.png"
qr_dp = "/tmp/qr_dp.png"
qr(GITHUB_URL, qr_gh)
qr(DEPLOY_URL, qr_dp)

qs = Inches(1.5)
qy = Inches(4.7)
# Centered under each card
qgx = sx5 + (cw5 - qs) // 2
qdx = dx5 + (cw5 - qs) // 2
s.shapes.add_picture(qr_gh, qgx, qy, qs, qs)
s.shapes.add_picture(qr_dp, qdx, qy, qs, qs)

# QR labels
tb = s.shapes.add_textbox(qgx - Inches(0.2), qy + qs + Inches(0.15), qs + Inches(0.4), Inches(0.35))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.text = "Scan: GitHub repo"; p.font.size = Pt(11); p.font.color.rgb = GRAY
p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER

tb = s.shapes.add_textbox(qdx - Inches(0.2), qy + qs + Inches(0.15), qs + Inches(0.4), Inches(0.35))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.text = "Scan: Live demo"; p.font.size = Pt(11); p.font.color.rgb = GRAY
p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER

# ======================== SAVE ========================
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "presentation.pptx")
prs.save(out)
print(f"Done → {out}")
